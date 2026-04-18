"""
파일 콘텐츠 추출 모듈

지원 형식:
  - PDF         : pdfplumber
  - DOCX        : python-docx
  - XLSX        : openpyxl
  - PPTX        : python-pptx
  - HWP 5.x     : pyhwp (hwp5)
  - 텍스트 계열 : 직접 읽기 (txt, csv, json, xml, md, py, js, ...)

SMB 파일은 임시 파일로 다운로드 후 파싱하고 삭제합니다.
"""

import io
import logging
import os
import tempfile
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

# 지원 확장자 → 추출 방법 분류
TEXT_EXTENSIONS = {
    ".txt", ".log", ".md", ".csv", ".json", ".xml",
    ".sql", ".yaml", ".yml", ".html", ".css",
    ".py", ".js", ".java", ".cpp", ".cs", ".ts",
}

OFFICE_EXTENSIONS = {
    ".pdf", ".docx", ".xlsx", ".pptx", ".doc", ".xls", ".ppt",
}

HWP_EXTENSIONS = {".hwp", ".hwpx"}

ALL_SUPPORTED = TEXT_EXTENSIONS | OFFICE_EXTENSIONS | HWP_EXTENSIONS

MAX_CONTENT_BYTES = 5 * 1024 * 1024   # SMB 다운로드 상한 5MB
MAX_CONTENT_CHARS = 50_000            # 인덱싱할 최대 글자 수


# ─────────────────────────────────────────────────────────
# 형식별 추출 함수 (로컬 파일 경로 받음)
# ─────────────────────────────────────────────────────────

def _extract_pdf(path: str) -> str:
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    texts.append(t)
        return "\n".join(texts)
    except Exception as e:
        logger.debug(f"PDF 추출 실패 ({path}): {e}")
        return ""


def _extract_docx(path: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        logger.debug(f"DOCX 추출 실패 ({path}): {e}")
        return ""


def _extract_xlsx(path: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    parts.append(row_text)
            if len("\n".join(parts)) > MAX_CONTENT_CHARS:
                break
        return "\n".join(parts)
    except Exception as e:
        logger.debug(f"XLSX 추출 실패 ({path}): {e}")
        return ""


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception as e:
        logger.debug(f"PPTX 추출 실패 ({path}): {e}")
        return ""


def _extract_hwp(path: str) -> str:
    """HWP 5.x 텍스트 추출

    1차: pyhwp CLI (hwp5txt) - 전체 본문
    2차: olefile PrvText 스트림 - 미리보기 텍스트 (폴백)
    """
    # 1차: pyhwp CLI 시도
    try:
        import subprocess, sys
        result = subprocess.run(
            [sys.executable, "-m", "hwp5.hwp5txt", path],
            capture_output=True, timeout=30
        )
        text = result.stdout.decode("utf-8", errors="replace").strip()
        if text:
            return text
    except Exception:
        pass

    # 2차: olefile PrvText (HWP OLE 컨테이너 미리보기 스트림)
    try:
        import olefile
        with olefile.OleFileIO(path) as ole:
            # PrvText: UTF-16-LE 인코딩 평문 미리보기
            if ole.exists("PrvText"):
                raw = ole.openstream("PrvText").read()
                return raw.decode("utf-16-le", errors="replace")
    except Exception as e:
        logger.debug(f"HWP OLE 추출 실패 ({path}): {e}")

    return ""


def _extract_text(path: str) -> str:
    """일반 텍스트 파일 읽기"""
    for enc in ("utf-8", "cp949", "euc-kr", "latin-1"):
        try:
            with open(path, "r", encoding=enc, errors="strict") as f:
                return f.read(MAX_CONTENT_CHARS)
        except (UnicodeDecodeError, LookupError):
            continue
        except Exception as e:
            logger.debug(f"텍스트 읽기 실패 ({path}): {e}")
            return ""
    return ""


# ─────────────────────────────────────────────────────────
# SMB 파일 → 임시 다운로드 → 추출
# ─────────────────────────────────────────────────────────

def _download_smb(smb_path: str, suffix: str) -> Optional[str]:
    """SMB 파일을 임시 파일로 다운로드. 실패 시 None 반환."""
    try:
        import smbclient
        with smbclient.open_file(smb_path, mode="rb") as src:
            data = src.read(MAX_CONTENT_BYTES)
        fd, tmp_path = tempfile.mkstemp(suffix=suffix)
        with os.fdopen(fd, "wb") as dst:
            dst.write(data)
        return tmp_path
    except Exception as e:
        logger.debug(f"SMB 다운로드 실패 ({smb_path}): {e}")
        return None


# ─────────────────────────────────────────────────────────
# 공개 API
# ─────────────────────────────────────────────────────────

def extract_content(file_path: str, file_size: int = 0) -> str:
    """
    파일 경로로부터 텍스트 콘텐츠 추출.

    Args:
        file_path: 로컬 경로 또는 SMB UNC 경로 (\\\\host\\share\\...)
        file_size: 파일 크기 (바이트). 0이면 크기 검사 건너뜀.

    Returns:
        추출된 텍스트 (최대 MAX_CONTENT_CHARS 글자). 실패 시 "".
    """
    ext = Path(file_path).suffix.lower()
    if ext not in ALL_SUPPORTED:
        return ""

    # 크기 제한 (다운로드 전 사전 차단)
    if file_size > MAX_CONTENT_BYTES:
        logger.debug(f"파일 크기 초과, 건너뜀: {file_path} ({file_size:,} bytes)")
        return ""

    is_smb = file_path.startswith("\\\\")
    tmp_path = None

    try:
        if is_smb:
            # 텍스트 파일은 직접 읽기 (임시 파일 불필요)
            if ext in TEXT_EXTENSIONS:
                try:
                    import smbclient
                    with smbclient.open_file(file_path, mode="rb") as f:
                        raw = f.read(MAX_CONTENT_BYTES)
                    for enc in ("utf-8", "cp949", "euc-kr", "latin-1"):
                        try:
                            return raw.decode(enc)[:MAX_CONTENT_CHARS]
                        except (UnicodeDecodeError, LookupError):
                            continue
                    return raw.decode("utf-8", errors="replace")[:MAX_CONTENT_CHARS]
                except Exception as e:
                    logger.debug(f"SMB 텍스트 읽기 실패 ({file_path}): {e}")
                    return ""
            else:
                # 바이너리 형식은 임시 파일로 다운로드
                tmp_path = _download_smb(file_path, suffix=ext)
                if not tmp_path:
                    return ""
                local_path = tmp_path
        else:
            local_path = file_path

        # 형식별 추출
        if ext == ".pdf":
            content = _extract_pdf(local_path)
        elif ext == ".docx":
            content = _extract_docx(local_path)
        elif ext == ".xlsx":
            content = _extract_xlsx(local_path)
        elif ext == ".pptx":
            content = _extract_pptx(local_path)
        elif ext in (".hwp", ".hwpx"):
            content = _extract_hwp(local_path)
        elif ext in TEXT_EXTENSIONS:
            content = _extract_text(local_path)
        else:
            content = ""

        return content[:MAX_CONTENT_CHARS]

    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.remove(tmp_path)
            except Exception:
                pass


def is_extractable(filename: str) -> bool:
    """해당 파일명의 확장자를 콘텐츠 추출할 수 있는지 확인"""
    return Path(filename).suffix.lower() in ALL_SUPPORTED
